import streamlit as st
import app.pages as pages
import app.utils as utils
import app.prompts as prompts
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO

# App title
pages.show_home()
pages.show_sidebar()

# Home page
st.title("✍️Style Writer")

# Content input
st.session_state.content = st.text_area(
    "Content Data:", st.session_state.content, 200
)

uploaded_files = st.file_uploader(
    "Upload Content Files:", 
    type=["pdf", "docx", "pptx"], 
    accept_multiple_files=True,
    help="Upload PDF, Word, or PowerPoint files"
)

# Extract text from uploaded files
extracted_text = ""
if uploaded_files:
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split('.')[-1].lower()
        
        if file_type == 'pdf':
            pdf_reader = PdfReader(BytesIO(uploaded_file.read()))
            for page in pdf_reader.pages:
                extracted_text += page.extract_text() + "\n"
                
        elif file_type == 'docx':
            doc = Document(BytesIO(uploaded_file.read()))
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    extracted_text += paragraph.text + "\n"
                    
        elif file_type == 'pptx':
            prs = Presentation(BytesIO(uploaded_file.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            extracted_text += shape.text + "\n"

# Combine text area and extracted content
content_all = st.session_state.content + "\n" + extracted_text.encode("ascii", errors="ignore").decode("ascii")

# Extracting the styles and creating combined display options
styles_data = utils.get_styles()
style_options = [item['name'] for item in styles_data]
selected_style = st.selectbox("Select a Style:", options=style_options, index=None)

# Assigning the selected style to the session state
if selected_style:
    # Find the matching style data
    filtered = next(
        (item for item in styles_data if str(item["name"]) == selected_style), None
    )
    
    if filtered:
        st.session_state.style = filtered["style"]
        st.session_state.example = filtered["example"]
        st.session_state.styleId = selected_style
        
st.session_state.style = st.text_area("✨Style", st.session_state.style)

# Show the example style
guidelines = st.session_state.locals.get("relevant_guidelines", {})
selected_guidelines = []

st.text("Editorial Style Guide (Select only the appropriate guidelines):")

# Create a checkbox for each guideline section
if guidelines:
    # Create two columns
    col1, col2 = st.columns(2)
    
    # Split guidelines into two halves
    guideline_items = list(guidelines.items())
    mid_point = len(guideline_items) // 2
    
    # First column
    with col1:
        for section_name, content in guideline_items[:mid_point]:
            default = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
            if st.checkbox(section_name, value=default, key=f"col1_{section_name}"):
                selected_guidelines.append(content)
    
    # Second column
    with col2:
        for section_name, content in guideline_items[mid_point:]:
            default = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
            if st.checkbox(section_name, value=default, key=f"col2_{section_name}"):
                selected_guidelines.append(content)
else:
    st.warning("No guidelines available in the local data.")

# Join all selected guidelines with newlines and store in session state
st.session_state.guidelines = "\n".join(selected_guidelines)

# Show the combined guidelines in a text area
st.text_area("✨Relevant Guidelines", st.session_state.guidelines, height=200)

if st.button(
    ":blue[Rewrite Content]",
    key="extract",
    disabled=content_all == ""
    or st.session_state.style == ""
    or st.session_state.example == "",
):
    with st.container(border=True):
        with st.spinner("Processing..."):
            output = prompts.rewrite_content(content_all, False)
            utils.save_output(output, content_all)
