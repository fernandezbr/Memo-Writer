import streamlit as st
import app.pages as pages
import app.utils as utils
import app.prompts as prompts
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO

# App title
pages.show_home()
pages.show_sidebar()

st.title("âœ¨Style Reader")

st.session_state.exampleText = st.text_area(
    "Reference Example:", st.session_state.exampleText, 200
)

uploaded_files = st.file_uploader(
    "Upload Example Files:", 
    type=["pdf", "docx", "pptx"], 
    accept_multiple_files=True,
    help="Upload PDF, Word, or PowerPoint files"
)
    
# Extract text from uploaded files
extracted_text = ""
if uploaded_files:
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split('.')[-1].lower()
        # Read the file content once
        file_content = uploaded_file.read()
        
        if file_type == 'pdf':
            pdf_reader = PdfReader(BytesIO(file_content))
            for page in pdf_reader.pages:
                extracted_text += page.extract_text() + "\n"
                
        elif file_type == 'docx':
            doc = Document(BytesIO(file_content))
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    extracted_text += paragraph.text + "\n"
                    
        elif file_type == 'pptx':
            prs = Presentation(BytesIO(file_content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            extracted_text += shape.text + "\n"

# Combine text area and extracted content
st.session_state.example = st.session_state.exampleText or extracted_text

if extracted_text:
    st.session_state.example = st.session_state.exampleText + "\n" + extracted_text

if st.button(
    ":blue[Extract Writing Style]",
    key="extract",
    disabled=st.session_state.example.strip() == "",
):
    with st.container(border=True):
        # Extract the writing style
        with st.spinner("Processing..."):
            st.session_state.style = prompts.extract_style(True)
            utils.save_style()
